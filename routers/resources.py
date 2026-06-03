"""Resource listing and detail routes with filtering, search, and pagination."""

import math
from fastapi import APIRouter, Depends, HTTPException, Query
from sqlalchemy.orm import Session
from sqlalchemy import func, or_
from database import get_db
from models import Resource, Category, Tag, resource_tags
from schemas import ResourceListItem, ResourceDetail, ResourceListResponse, TagOut

router = APIRouter(prefix="/api/v1/resources", tags=["resources"])


def _resource_to_listitem(resource: Resource) -> ResourceListItem:
    return ResourceListItem(
        id=resource.id,
        title=resource.title,
        slug=resource.slug,
        description=resource.description or "",
        category_slug=resource.category.slug if resource.category else "",
        category_name=resource.category.name if resource.category else "",
        file_type=resource.file_type or "",
        file_size=resource.file_size or "",
        author=resource.author or "",
        difficulty=resource.difficulty or 1,
        view_count=resource.view_count or 0,
        download_count=resource.download_count or 0,
        is_featured=resource.is_featured or False,
        tags=[TagOut(id=t.id, name=t.name, slug=t.slug, resource_count=0)
              for t in (resource.tags or [])],
        created_at=resource.created_at,
    )


def _resource_to_detail(resource: Resource) -> ResourceDetail:
    return ResourceDetail(
        id=resource.id,
        title=resource.title,
        slug=resource.slug,
        description=resource.description or "",
        content=resource.content or "",
        category_slug=resource.category.slug if resource.category else "",
        category_name=resource.category.name if resource.category else "",
        file_url=resource.file_url or "",
        file_type=resource.file_type or "",
        file_size=resource.file_size or "",
        author=resource.author or "",
        source=resource.source or "",
        difficulty=resource.difficulty or 1,
        view_count=resource.view_count or 0,
        download_count=resource.download_count or 0,
        is_featured=resource.is_featured or False,
        tags=[TagOut(id=t.id, name=t.name, slug=t.slug, resource_count=0)
              for t in (resource.tags or [])],
        created_at=resource.created_at,
        updated_at=resource.updated_at,
    )


@router.get("", response_model=ResourceListResponse)
def list_resources(
    category: str = Query(None, description="Filter by category slug"),
    tag: str = Query(None, description="Filter by tag slug"),
    search: str = Query(None, description="Search in title and description"),
    difficulty: int = Query(None, ge=1, le=5, description="Filter by difficulty 1-5"),
    featured: bool = Query(None, description="Show featured only"),
    sort: str = Query("newest", description="Sort: newest, popular, featured"),
    page: int = Query(1, ge=1, description="Page number"),
    page_size: int = Query(20, ge=1, le=100, description="Items per page"),
    db: Session = Depends(get_db),
):
    query = db.query(Resource).filter(Resource.status == "published")

    # ── Filters ────────────────────────────────────
    if category:
        cat_obj = db.query(Category).filter(Category.slug == category).first()
        if cat_obj:
            query = query.filter(Resource.category_id == cat_obj.id)

    if tag:
        tag_obj = db.query(Tag).filter(Tag.slug == tag).first()
        if tag_obj:
            query = query.join(resource_tags).filter(resource_tags.c.tag_id == tag_obj.id)

    if search:
        like = f"%{search}%"
        query = query.filter(
            or_(Resource.title.ilike(like), Resource.description.ilike(like))
        )

    if difficulty is not None:
        query = query.filter(Resource.difficulty == difficulty)

    if featured:
        query = query.filter(Resource.is_featured == True)

    # ── Sorting ────────────────────────────────────
    if sort == "popular":
        query = query.order_by(Resource.view_count.desc())
    elif sort == "featured":
        query = query.order_by(Resource.is_featured.desc(), Resource.created_at.desc())
    else:  # newest
        query = query.order_by(Resource.created_at.desc())

    # ── Pagination ─────────────────────────────────
    total = query.count()
    total_pages = max(1, math.ceil(total / page_size))
    offset = (page - 1) * page_size
    resources = query.offset(offset).limit(page_size).all()

    items = [_resource_to_listitem(r) for r in resources]

    return ResourceListResponse(
        items=items,
        total=total,
        page=page,
        page_size=page_size,
        total_pages=total_pages,
    )


@router.get("/featured", response_model=list[ResourceListItem])
def list_featured(db: Session = Depends(get_db)):
    resources = (
        db.query(Resource)
        .filter(Resource.status == "published", Resource.is_featured == True)
        .order_by(Resource.created_at.desc())
        .limit(6)
        .all()
    )
    return [_resource_to_listitem(r) for r in resources]


@router.get("/{resource_id}", response_model=ResourceDetail)
def get_resource(resource_id: int, db: Session = Depends(get_db)):
    resource = db.query(Resource).filter(Resource.id == resource_id).first()
    if not resource:
        raise HTTPException(status_code=404, detail="Resource not found")

    # Increment view count
    resource.view_count = (resource.view_count or 0) + 1
    db.commit()

    return _resource_to_detail(resource)


# ── File listing & upload ──────────────────────────────

import os
from urllib.parse import quote
from security import secure_filename, upload_limiter
from pathlib import Path
from fastapi import UploadFile, File as FileParam, Request

UPLOADS_DIR = Path(__file__).resolve().parent.parent / "static" / "uploads"
MAX_UPLOAD_SIZE = 100 * 1024 * 1024  # 100 MB


@router.get("/{resource_id}/files")
def list_resource_files(resource_id: int):
    """List all files in a resource's upload folder."""
    folder = UPLOADS_DIR / str(resource_id)
    if not folder.exists():
        return {"resource_id": resource_id, "files": [], "folder_url": f"/uploads/{resource_id}/"}

    files = []
    for f in sorted(folder.rglob("*")):
        if not f.is_file():
            continue
        rel = f.relative_to(folder)
        if f.is_file():
            stat = f.stat()
            size = stat.st_size
            if size < 1024:
                size_str = f"{size}B"
            elif size < 1024 * 1024:
                size_str = f"{size / 1024:.1f}KB"
            else:
                size_str = f"{size / (1024 * 1024):.1f}MB"

            ext = f.suffix.lower()
            preview_type = "none"
            if ext in (".pdf",):
                preview_type = "pdf"
            elif ext in (".jpg", ".jpeg", ".png", ".gif", ".webp", ".svg"):
                preview_type = "image"
            elif ext in (".txt", ".md", ".py", ".js", ".html", ".css", ".json", ".xml", ".csv"):
                preview_type = "text"
            elif ext in (".pptx", ".ppt"):
                preview_type = "pptx"
            elif ext in (".docx", ".doc"):
                preview_type = "docx"
            elif ext in (".mp4", ".webm", ".mov"):
                preview_type = "video"
            elif ext in (".mp3", ".wav", ".ogg"):
                preview_type = "audio"

            files.append({
                "name": str(rel).replace("\\", "/"),
                "size": size_str,
                "size_bytes": size,
                "url": f"/uploads/{resource_id}/" + "/".join(quote(p) for p in str(rel).replace("\\", "/").split("/")),
                "preview_type": preview_type,
                "extension": ext,
                "modified": stat.st_mtime,
            })

    return {
        "resource_id": resource_id,
        "folder_url": f"/uploads/{resource_id}/",
        "files": files,
    }


@router.post("/{resource_id}/upload")
async def upload_resource_file(request: Request, resource_id: int, file: UploadFile = FileParam(...)):
    """Upload a file to a resource's folder (max 100MB, rate limited)."""
    # Rate limit
    upload_limiter.limit(request)

    # Sanitize filename to prevent path injection
    safe_name = secure_filename(file.filename)

    folder = UPLOADS_DIR / str(resource_id)
    folder.mkdir(parents=True, exist_ok=True)

    file_path = folder / safe_name
    raw = await file.read()

    # Size limit check
    if len(raw) > MAX_UPLOAD_SIZE:
        raise HTTPException(status_code=413, detail=f"File too large. Max {MAX_UPLOAD_SIZE // (1024*1024)}MB")

    file_path.write_bytes(raw)

    ext = file_path.suffix.lower()
    preview_type = "none"
    if ext in (".pdf",): preview_type = "pdf"
    elif ext in (".jpg", ".jpeg", ".png", ".gif", ".webp", ".svg"): preview_type = "image"
    elif ext in (".txt", ".md", ".py", ".js", ".html", ".css", ".json", ".xml", ".csv"): preview_type = "text"
    elif ext in (".pptx", ".ppt"): preview_type = "pptx"
    elif ext in (".docx", ".doc"): preview_type = "docx"

    return {
        "resource_id": resource_id,
        "name": safe_name,
        "size": len(raw),
        "url": f"/uploads/{resource_id}/{safe_name}",
        "preview_type": preview_type,
    }


# ── PPTX / Office Preview ────────────────────────────────

from fastapi.responses import HTMLResponse as HTMLResp

@router.get("/{resource_id}/preview/{filename:path}", response_class=HTMLResp)
def preview_office_file(resource_id: int, filename: str):
    """Convert PPTX/DOCX to HTML slideshow for inline preview."""
    folder = UPLOADS_DIR / str(resource_id)
    file_path = folder / filename
    if not file_path.exists():
        raise HTTPException(status_code=404, detail="File not found")

    ext = file_path.suffix.lower()
    if ext == ".pptx":
        try:
            from pptx import Presentation
            from pptx.util import Inches
            from pptx.enum.shapes import MSO_SHAPE_TYPE
            import base64

            prs = Presentation(str(file_path))
            sw = prs.slide_width or Inches(10)
            sh = prs.slide_height or Inches(7.5)
            ratio = sh / sw  # preserve aspect ratio

            def pct(val, base):
                """Safe percentage from EMU."""
                try:
                    return (val / base) * 100 if val is not None else 0
                except Exception:
                    return 0

            def make_style(left, top, width, height, auto_height=False):
                l = pct(left, sw); t = pct(top, sh)
                w = pct(width, sw); h = pct(height, sh)
                h_rule = f"min-height:{h:.1f}%;" if auto_height else f"height:{h:.1f}%;"
                return f"position:absolute;left:{l:.1f}%;top:{t:.1f}%;width:{w:.1f}%;{h_rule}"

            def render_textframe(tf):
                """Render a text frame with runs formatting and alignment."""
                parts = []
                for para in tf.paragraphs:
                    # Paragraph alignment
                    align_css = ""
                    try:
                        from pptx.enum.text import PP_ALIGN
                        al = para.alignment
                        if al == PP_ALIGN.CENTER: align_css = "text-align:center;"
                        elif al == PP_ALIGN.RIGHT: align_css = "text-align:right;"
                        elif al == PP_ALIGN.JUSTIFY: align_css = "text-align:justify;"
                    except Exception:
                        pass

                    # Bullet/indent
                    bullet = ""
                    indent_css = ""
                    try:
                        lvl = para.level or 0
                        if lvl > 0:
                            indent_css = f"padding-left:{lvl*1.5}em;"
                    except Exception:
                        pass

                    run_parts = []
                    for run in para.runs:
                        try:
                            fz = run.font.size.pt if run.font.size else None
                        except Exception:
                            fz = None
                        if fz:
                            css = f"font-size:clamp(8px,{fz/7.2:.1f}vw,{fz*1.5}px);"
                        else:
                            css = "font-size:clamp(8px,1.5vw,18px);"
                        if run.font.bold: css += "font-weight:bold;"
                        if run.font.italic: css += "font-style:italic;"
                        if run.font.underline: css += "text-decoration:underline;"
                        has_color = False
                        try:
                            if run.font.color and run.font.color.rgb:
                                raw = str(run.font.color.rgb)
                                # python-pptx sometimes returns hex, sometimes decimal
                                if len(raw) <= 6 and all(c in '0123456789ABCDEFabcdef' for c in raw):
                                    hex_c = raw.upper().zfill(6)
                                    rgb_dec = int(hex_c, 16)
                                else:
                                    rgb_dec = int(raw)
                                    hex_c = f"{rgb_dec:06X}"
                                r_val = (rgb_dec >> 16) & 0xFF
                                g_val = (rgb_dec >> 8) & 0xFF
                                b_val = rgb_dec & 0xFF
                                brightness = (r_val * 299 + g_val * 587 + b_val * 114) / 1000
                                if brightness < 200:
                                    css += f"color:#{hex_c};"
                                    has_color = True
                        except Exception:
                            pass
                        if not has_color:
                            css += "color:#1a1a18;"
                        t = run.text.replace("&","&amp;").replace("<","&lt;").replace(">","&gt;").replace("\n","<br>")
                        run_parts.append(f'<span style="{css}">{t}</span>')
                    if run_parts:
                        parts.append(f'<p style="margin:0.1em 0;line-height:1.35;{align_css}{indent_css}">{"".join(run_parts)}</p>')
                    else:
                        t = para.text.replace("&","&amp;").replace("<","&lt;").replace(">","&gt;")
                        if t.strip():
                            parts.append(f'<p style="margin:0.1em 0;line-height:1.35;color:#1a1a18;{align_css}{indent_css}">{t}</p>')
                        else:
                            parts.append(f'<p style="margin:0.3em 0;line-height:1">&#8203;</p>')
                return "".join(parts)

            def render_shape(shape, output_list):
                """Render a single shape (recursive for groups)."""
                try:
                    st = shape.shape_type
                except Exception:
                    st = None

                # Group shapes - recurse
                if st == MSO_SHAPE_TYPE.GROUP:
                    try:
                        for child in shape.shapes:
                            render_shape(child, output_list)
                    except Exception:
                        pass
                    return

                style_fixed = make_style(shape.left, shape.top, shape.width, shape.height)
                style_auto = make_style(shape.left, shape.top, shape.width, shape.height, auto_height=True)

                # Picture
                if st == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        img = shape.image
                        b64 = base64.b64encode(img.blob).decode()
                        ext_img = img.ext or "png"
                        output_list.append(
                            f'<div style="{style_fixed};display:flex;align-items:center;justify-content:center;overflow:hidden">'
                            f'<img src="data:image/{ext_img};base64,{b64}" style="max-width:100%;max-height:100%;object-fit:contain"></div>'
                        )
                    except Exception:
                        pass
                    return

                # Table
                try:
                    if shape.has_table:
                        rows_html = []
                        for row in shape.table.rows:
                            cells = []
                            for cell in row.cells:
                                ct = cell.text.replace("&","&amp;").replace("<","&lt;").replace(">","&gt;")
                                cells.append(f'<td style="border:1px solid #ccc;padding:3px 6px;font-size:10px;vertical-align:top">{ct}</td>')
                            rows_html.append(f'<tr>{"".join(cells)}</tr>')
                        output_list.append(
                            f'<div style="{style_fixed};overflow:auto">'
                            f'<table style="border-collapse:collapse;width:100%;font-size:10px">{"".join(rows_html)}</table></div>'
                        )
                        return
                except Exception:
                    pass

                # Text frame (auto-shape, text box, placeholder, title, etc.)
                try:
                    if shape.has_text_frame:
                        text_html = render_textframe(shape.text_frame)
                        if text_html.strip():
                            # Shape fill color
                            shape_bg = ""
                            try:
                                fill = shape.fill
                                if fill.type is not None:
                                    try:
                                        fc = fill.fore_color
                                        if fc and fc.rgb:
                                            raw = str(fc.rgb)
                                            if len(raw) <= 6 and all(c in '0123456789ABCDEFabcdef' for c in raw):
                                                shape_bg = f"background:#{raw.upper().zfill(6)};"
                                            else:
                                                shape_bg = f"background:#{int(raw):06X};"
                                    except Exception:
                                        pass
                            except Exception:
                                pass
                            output_list.append(
                                f'<div style="{style_auto};overflow:visible;padding:6px 8px;word-wrap:break-word;word-break:break-word;{shape_bg}">'
                                f'{text_html}</div>'
                            )
                        return
                except Exception:
                    pass

                # Fallback: try to get any text
                try:
                    if hasattr(shape, 'text') and shape.text.strip():
                        t = shape.text.replace("&","&amp;").replace("<","&lt;").replace(">","&gt;")
                        output_list.append(
                            f'<div style="{style_auto};overflow:visible;padding:4px 6px;font-size:clamp(8px,1.2vw,16px);word-wrap:break-word">'
                            f'<p style="margin:0;line-height:1.3">{t}</p></div>'
                        )
                except Exception:
                    pass

            slides_html = []
            for si, slide in enumerate(prs.slides):
                # Slide background
                bg_color = "#f5f5f5"
                try:
                    bg = slide.background
                    fill = bg.fill
                    if fill.type is not None:
                        try:
                            fc = fill.fore_color
                            if fc and fc.rgb:
                                raw = str(fc.rgb)
                                if len(raw) <= 6 and all(c in '0123456789ABCDEFabcdef' for c in raw):
                                    bg_color = f"#{raw.upper().zfill(6)}"
                                else:
                                    bg_color = f"#{int(raw):06X}"
                        except Exception:
                            pass
                except Exception:
                    pass

                shapes_parts = []
                for shape in slide.shapes:
                    render_shape(shape, shapes_parts)

                slides_html.append(
                    f'<div class="slide" data-slide="{si}" style="display:{"block" if si==0 else "none"};'
                    f'position:relative;width:100%;padding-top:{ratio*100:.1f}%;'
                    f'background:{bg_color};border:1px solid #ddd;overflow:visible;margin-bottom:20px">'
                    f'<div style="position:absolute;inset:0;overflow:visible">{"".join(shapes_parts)}</div></div>'
                )

            controls = ""
            if len(slides_html) > 1:
                controls = f'''
<div style="display:flex;align-items:center;justify-content:center;gap:8px;padding:8px;background:#f5f5f5;border:1px solid #e0e0e0;border-top:none">
<button onclick="prevSlide()" style="padding:4px 12px;border:1px solid #ccc;border-radius:4px;background:#fff;cursor:pointer;font-size:13px">&#8592; 上一页</button>
<span id="slideNum" style="font-size:13px;min-width:60px;text-align:center">1 / {len(slides_html)}</span>
<button onclick="nextSlide()" style="padding:4px 12px;border:1px solid #ccc;border-radius:4px;background:#fff;cursor:pointer;font-size:13px">下一页 &#8594;</button>
</div>'''

            html = f'''<!DOCTYPE html>
<html><head><meta charset="UTF-8"><meta name="viewport" content="width=device-width,initial-scale=1">
<title>{filename}</title>
<style>
*{{margin:0;padding:0;box-sizing:border-box}}
body{{font-family:'Instrument Sans','Microsoft YaHei',sans-serif;background:#e8e8e8;padding:20px}}
.container{{max-width:960px;margin:0 auto}}
h2{{font-size:14px;color:#666;margin-bottom:12px;text-align:center}}
</style></head>
<body>
<div class="container">
<h2>{filename} （{len(prs.slides)} 页）</h2>
<div style="background:#fff;box-shadow:0 2px 12px rgba(0,0,0,.1)">{"".join(slides_html)}</div>
{controls}
</div>
<script>
var total={len(slides_html)};var cur=0;
function showSlide(n){{cur=(n+total)%total;var slides=document.querySelectorAll(".slide");slides.forEach(function(s,i){{s.style.display=i===cur?"block":"none"}});document.getElementById("slideNum").textContent=(cur+1)+" / "+total}}
function prevSlide(){{showSlide(cur-1)}}
function nextSlide(){{showSlide(cur+1)}}
document.addEventListener("keydown",function(e){{if(e.key==="ArrowLeft")prevSlide();if(e.key==="ArrowRight")nextSlide()}})
</script></body></html>'''
            return HTMLResp(content=html)

        except ImportError:
            return HTMLResp(content="<p>python-pptx 未安装，无法预览 PPT。请运行: pip install python-pptx</p>")
        except Exception as e:
            return HTMLResp(content=f"<p>PPT 解析失败: {e}</p>")

    return HTMLResp(content=f"<p>不支持预览此文件格式: {ext}</p>")
